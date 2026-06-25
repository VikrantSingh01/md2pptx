"""Render a :class:`~md2pptx.model.Deck` to a ``.pptx`` using ``python-pptx``.

The renderer drives every slide off the blank layout and positions shapes by
hand. That is more work than relying on placeholders, but it gives pixel-level
control over typography, spacing and accent rules - which is what separates a
"markdown dump" from a leadership-grade deck.
"""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from .inline import parse_inline
from .model import Bullet, Column, Deck, Slide, SlideKind
from .themes import Theme, get_theme

# 16:9 canvas.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str)


class Renderer:
    def __init__(self, deck: Deck):
        self.deck = deck
        self.theme: Theme = get_theme(deck.theme)
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self._blank = self.prs.slide_layouts[6]

    # ------------------------------------------------------------------ #
    def render(self) -> Presentation:
        for index, slide in enumerate(self.deck.slides):
            self._render_slide(slide, index)
        return self.prs

    def save(self, path: str) -> None:
        self.render().save(path)

    # ------------------------------------------------------------------ #
    def _render_slide(self, slide: Slide, index: int) -> None:
        s = self.prs.slides.add_slide(self._blank)
        dark = slide.kind in (SlideKind.TITLE, SlideKind.SECTION, SlideKind.CLOSING)
        self._paint_background(s, dark)

        if slide.kind == SlideKind.TITLE:
            self._title_slide(s, slide)
        elif slide.kind == SlideKind.SECTION:
            self._section_slide(s, slide, index)
        elif slide.kind == SlideKind.CLOSING:
            self._closing_slide(s, slide)
        elif slide.kind == SlideKind.QUOTE:
            self._quote_slide(s, slide)
        elif slide.kind == SlideKind.IMAGE:
            self._image_slide(s, slide)
        elif slide.kind == SlideKind.TWO_COLUMN:
            self._two_column_slide(s, slide)
        else:
            self._content_slide(s, slide)

        if not dark:
            self._footer(s, index)
        if slide.notes:
            s.notes_slide.notes_text_frame.text = slide.notes

    # ------------------------------------------------------------------ #
    # Backgrounds & chrome
    # ------------------------------------------------------------------ #
    def _paint_background(self, s, dark: bool) -> None:
        pal = self.theme.palette
        color = pal.primary if dark else pal.background
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        rect.fill.solid()
        rect.fill.fore_color.rgb = _rgb(color)
        rect.line.fill.background()
        rect.shadow.inherit = False
        # Send to back.
        s.shapes._spTree.remove(rect._element)
        s.shapes._spTree.insert(2, rect._element)

    def _accent_bar(self, s, top: Emu, height=None, left=None, width=None) -> None:
        height = Pt(4) if height is None else height
        left = self.m if left is None else left
        width = Inches(0.9) if width is None else width
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(self.theme.palette.accent)
        bar.line.fill.background()
        bar.shadow.inherit = False

    def _footer(self, s, index: int) -> None:
        pal = self.theme.palette
        text = self.deck.footer or self.deck.title
        if text:
            tb = self._textbox(s, self.m, Inches(7.02), Inches(9.5), Inches(0.35))
            self._set_text(tb, text, self.theme.small_size, pal.secondary)
        # Page number.
        num = self._textbox(s, Inches(12.0), Inches(7.02), Inches(0.9), Inches(0.35))
        self._set_text(num, str(index + 1), self.theme.small_size, pal.secondary,
                       align=PP_ALIGN.RIGHT)

    # ------------------------------------------------------------------ #
    # Layouts
    # ------------------------------------------------------------------ #
    def _title_slide(self, s, slide: Slide) -> None:
        pal = self.theme.palette
        self._accent_bar(s, Inches(2.55), height=Pt(5), width=Inches(1.4))
        title = self._textbox(s, self.m, Inches(2.75), Inches(11.3), Inches(2.0))
        self._set_text(title, slide.title or self.deck.title, self.theme.title_size + 12,
                       pal.text_inverse, bold=True)
        sub = slide.subtitle or self.deck.subtitle
        if sub:
            sb = self._textbox(s, self.m, Inches(4.7), Inches(11.0), Inches(0.9))
            self._set_text(sb, sub, self.theme.heading_size - 6, pal.accent_soft)
        meta = " · ".join(x for x in (self.deck.author, self.deck.date) if x)
        if meta:
            mb = self._textbox(s, self.m, Inches(6.4), Inches(11.0), Inches(0.5))
            self._set_text(mb, meta, self.theme.body_size, pal.secondary)

    def _section_slide(self, s, slide: Slide, index: int) -> None:
        pal = self.theme.palette
        eyebrow = slide.eyebrow or "Section"
        eb = self._textbox(s, self.m, Inches(2.6), Inches(11.0), Inches(0.6))
        self._set_text(eb, eyebrow.upper(), self.theme.body_size, pal.accent,
                       bold=True, spacing=2.0)
        self._accent_bar(s, Inches(3.3), height=Pt(4), width=Inches(1.1))
        title = self._textbox(s, self.m, Inches(3.5), Inches(11.3), Inches(2.2))
        self._set_text(title, slide.title, self.theme.section_size + 6,
                       pal.text_inverse, bold=True)

    def _closing_slide(self, s, slide: Slide) -> None:
        pal = self.theme.palette
        box = self._textbox(s, self.m, Inches(3.0), Inches(11.3), Inches(1.6))
        self._set_text(box, slide.title or "Thank you", self.theme.title_size,
                       pal.text_inverse, bold=True)
        sub = slide.subtitle or "; ".join(p for p in slide.paragraphs)
        if sub:
            sb = self._textbox(s, self.m, Inches(4.7), Inches(11.0), Inches(1.4))
            self._set_text(sb, sub, self.theme.heading_size - 8, pal.accent_soft)

    def _quote_slide(self, s, slide: Slide) -> None:
        pal = self.theme.palette
        q = slide.quote
        mark = self._textbox(s, self.m, Inches(1.4), Inches(2.0), Inches(1.6))
        self._set_text(mark, "\u201c", 130, pal.accent, bold=True)
        body = self._textbox(s, self.m, Inches(2.6), Inches(11.0), Inches(3.0))
        self._set_text(body, q.text, self.theme.heading_size, pal.primary,
                       italic=True, line_spacing=1.15)
        if q.attribution:
            attr = self._textbox(s, self.m, Inches(5.7), Inches(11.0), Inches(0.6))
            self._set_text(attr, f"- {q.attribution}", self.theme.body_size,
                           pal.secondary, bold=True)

    def _image_slide(self, s, slide: Slide) -> None:
        if slide.title:
            self._heading(s, slide)
        self._place_image(s, slide.image, Inches(1.4), Inches(2.0),
                          max_w=Inches(10.5), max_h=Inches(4.6))

    def _content_slide(self, s, slide: Slide) -> None:
        self._heading(s, slide)
        top = Inches(2.05)
        if slide.paragraphs and not slide.bullets:
            tb = self._textbox(s, self.m, top, Inches(11.5), Inches(4.5))
            self._fill_paragraphs(tb, slide.paragraphs)
        elif slide.bullets:
            tb = self._textbox(s, self.m, top, Inches(11.5), Inches(4.6))
            self._fill_bullets(tb, slide.bullets)
        if slide.table:
            self._place_table(s, slide, top if not slide.bullets else Inches(4.4))
        if slide.code:
            self._place_code(s, slide.code, top if not slide.bullets else Inches(4.2))
        if slide.image and not (slide.bullets or slide.table or slide.code):
            self._place_image(s, slide.image, Inches(1.6), Inches(2.2),
                              max_w=Inches(10.0), max_h=Inches(4.4))

    def _two_column_slide(self, s, slide: Slide) -> None:
        self._heading(s, slide)
        cols = slide.columns[:2]
        gutter = Inches(0.5)
        usable = SLIDE_W - self.m * 2 - gutter
        col_w = Emu(int(usable / 2))
        lefts = [self.m, Emu(int(self.m + col_w + gutter))]
        for col, left in zip(cols, lefts):
            self._render_column(s, col, left, col_w)

    def _render_column(self, s, col: Column, left, width) -> None:
        top = Inches(2.1)
        if col.image:
            self._place_image(s, col.image, left, top, max_w=width, max_h=Inches(4.4))
            return
        tb = self._textbox(s, left, top, width, Inches(4.6))
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for para in col.paragraphs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            self._apply_runs(p, para, self.theme.body_size, self.theme.palette.text)
            p.space_after = Pt(8)
        for b in col.bullets:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            self._style_bullet(p, b)

    # ------------------------------------------------------------------ #
    # Shared building blocks
    # ------------------------------------------------------------------ #
    def _heading(self, s, slide: Slide) -> None:
        pal = self.theme.palette
        if slide.eyebrow:
            eb = self._textbox(s, self.m, Inches(0.55), Inches(11.0), Inches(0.4))
            self._set_text(eb, slide.eyebrow.upper(), self.theme.small_size + 1,
                           pal.accent, bold=True, spacing=2.0)
            title_top = Inches(0.95)
        else:
            title_top = Inches(0.7)
        title = self._textbox(s, self.m, title_top, Inches(11.5), Inches(1.0))
        self._set_text(title, slide.title, self.theme.heading_size, pal.primary, bold=True)
        self._accent_bar(s, Inches(1.85), height=Pt(3), width=Inches(0.9))

    def _fill_bullets(self, tb, bullets: list[Bullet]) -> None:
        tf = tb.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            self._style_bullet(p, b)

    def _style_bullet(self, p, b: Bullet) -> None:
        pal = self.theme.palette
        p.level = b.level
        size = self.theme.body_size - (1 if b.level else 0)
        self._add_run(p, "▪  " if b.level == 0 else "–  ", size, pal.accent, bold=True)
        self._apply_runs(p, b.text, size, pal.text, append=True)
        p.space_after = Pt(7)
        p.line_spacing = 1.08

    def _fill_paragraphs(self, tb, paragraphs: list[str]) -> None:
        tf = tb.text_frame
        tf.word_wrap = True
        for i, para in enumerate(paragraphs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            self._apply_runs(p, para, self.theme.body_size, self.theme.palette.text)
            p.space_after = Pt(10)
            p.line_spacing = 1.12

    def _place_table(self, s, slide: Slide, top) -> None:
        t = slide.table
        rows = len(t.rows) + 1
        cols = len(t.headers)
        width = SLIDE_W - self.m * 2
        height = Inches(min(0.5 * rows, 4.5))
        shape = s.shapes.add_table(rows, cols, self.m, top, width, height)
        table = shape.table
        pal = self.theme.palette
        for c, head in enumerate(t.headers):
            cell = table.cell(0, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(pal.primary)
            self._cell_text(cell, head, pal.text_inverse, bold=True)
        for r, row in enumerate(t.rows, start=1):
            for c in range(cols):
                cell = table.cell(r, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(pal.surface if r % 2 else pal.background)
                value = row[c] if c < len(row) else ""
                self._cell_text(cell, value, pal.text)

    def _cell_text(self, cell, text, color, bold=False) -> None:
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_top = Pt(4)
        cell.margin_bottom = Pt(4)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        self._apply_runs(p, text, self.theme.small_size + 2, color, bold=bold)

    def _place_code(self, s, code, top) -> None:
        bottom_limit = Inches(6.85)  # keep clear of the footer
        height = Emu(max(int(bottom_limit - top), int(Inches(1.2))))
        panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, self.m, top,
                                   SLIDE_W - self.m * 2, height)
        panel.fill.solid()
        panel.fill.fore_color.rgb = _rgb("0B1220")
        panel.line.fill.background()
        panel.shadow.inherit = False
        tf = panel.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Inches(0.35)
        tf.margin_right = Inches(0.35)
        tf.margin_top = Inches(0.25)
        tf.margin_bottom = Inches(0.25)
        for i, line in enumerate(code.code.splitlines() or [""]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.05
            run = p.add_run()
            run.text = line or " "
            run.font.name = self.theme.fonts.mono
            run.font.size = Pt(self.theme.small_size + 2)
            run.font.color.rgb = _rgb("E6EDF5")

    def _place_image(self, s, image, left, top, max_w, max_h) -> None:
        from os.path import exists

        if not image or (not image.path.startswith("http") and not exists(image.path)):
            self._missing_image(s, image, left, top, max_w, max_h)
            return
        try:
            pic = s.shapes.add_picture(image.path, left, top)
        except Exception:
            self._missing_image(s, image, left, top, max_w, max_h)
            return
        scale = min(max_w / pic.width, max_h / pic.height, 1.0)
        pic.width = Emu(int(pic.width * scale))
        pic.height = Emu(int(pic.height * scale))
        pic.left = Emu(int((SLIDE_W - pic.width) / 2))

    def _missing_image(self, s, image, left, top, max_w, max_h) -> None:
        pal = self.theme.palette
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, max_w, max_h)
        box.fill.solid()
        box.fill.fore_color.rgb = _rgb(pal.surface)
        box.line.color.rgb = _rgb(pal.secondary)
        box.shadow.inherit = False
        label = (image.alt if image and image.alt else "image")
        self._set_text(box, f"[ {label} ]", self.theme.body_size, pal.secondary,
                       align=PP_ALIGN.CENTER)

    # ------------------------------------------------------------------ #
    # Low-level text helpers
    # ------------------------------------------------------------------ #
    @property
    def m(self) -> Emu:
        return Inches(self.theme.margin)

    def _textbox(self, s, left, top, width, height):
        return s.shapes.add_textbox(left, top, width, height)

    def _set_text(self, shape, text, size, color, *, bold=False, italic=False,
                  align=PP_ALIGN.LEFT, spacing=None, line_spacing=None) -> None:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        f = run.font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.name = self.theme.fonts.heading if bold else self.theme.fonts.body
        f.color.rgb = _rgb(color)
        if spacing:
            self._letter_spacing(run, spacing)

    def _apply_runs(self, p, text, size, color, *, bold=False, append=False) -> None:
        for run in parse_inline(text):
            r = p.add_run()
            r.text = run.text
            f = r.font
            f.size = Pt(size - 1 if run.code else size)
            f.bold = bold or run.bold
            f.italic = run.italic
            f.name = self.theme.fonts.mono if run.code else self.theme.fonts.body
            f.color.rgb = _rgb(self.theme.palette.accent if run.code else color)

    def _add_run(self, p, text, size, color, *, bold=False):
        r = p.add_run()
        r.text = text
        f = r.font
        f.size = Pt(size)
        f.bold = bold
        f.name = self.theme.fonts.body
        f.color.rgb = _rgb(color)
        return r

    def _letter_spacing(self, run, points: float) -> None:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(points * 100)))


def render_deck(deck: Deck, output_path: str) -> None:
    Renderer(deck).save(output_path)
