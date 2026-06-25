"""End-to-end rendering tests: Markdown → real .pptx on disk."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from md2pptx.builder import build

EXAMPLES = Path(__file__).resolve().parent.parent / "examples"


def test_build_demo_deck(tmp_path):
    out = tmp_path / "demo.pptx"
    deck = build([EXAMPLES / "leadership-demo.md"], out)

    assert out.exists() and out.stat().st_size > 0
    prs = Presentation(str(out))
    # Every parsed slide should appear in the file.
    assert len(prs.slides) == len(deck.slides) >= 10
    # 16:9 canvas.
    assert prs.slide_width == Inches(13.333)


def test_build_multifile_concatenates(tmp_path):
    out = tmp_path / "wbr.pptx"
    files = sorted((EXAMPLES / "multi").glob("*.md"))
    deck = build(files, out)
    prs = Presentation(str(out))
    assert len(prs.slides) == len(deck.slides)
    # First file's front matter sets deck metadata.
    assert deck.title == "Weekly Business Review"
    assert deck.theme == "midnight"


def test_theme_override(tmp_path):
    out = tmp_path / "themed.pptx"
    deck = build([EXAMPLES / "leadership-demo.md"], out, theme="slate")
    assert deck.theme == "slate"
    assert out.exists()


def test_unknown_theme_override_raises(tmp_path):
    with pytest.raises(KeyError):
        build([EXAMPLES / "leadership-demo.md"], tmp_path / "x.pptx", theme="nope")


def test_notes_are_written(tmp_path):
    out = tmp_path / "notes.pptx"
    build([EXAMPLES / "leadership-demo.md"], out)
    prs = Presentation(str(out))
    all_notes = " ".join(
        s.notes_slide.notes_text_frame.text
        for s in prs.slides
        if s.has_notes_slide
    )
    assert "headline numbers" in all_notes


def test_empty_inputs_raises(tmp_path):
    with pytest.raises(ValueError):
        build([], tmp_path / "x.pptx")
